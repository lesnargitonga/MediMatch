from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
OUT = ROOT / "MediMatch_Conference_Deck.pptx"

W = Inches(13.333)
H = Inches(7.5)

BG = RGBColor(7, 12, 24)
PANEL = RGBColor(18, 29, 49)
PANEL_2 = RGBColor(26, 38, 59)
CREAM = RGBColor(248, 243, 231)
MUTED = RGBColor(174, 186, 204)
GOLD = RGBColor(246, 190, 76)
ORANGE = RGBColor(255, 119, 68)
MINT = RGBColor(67, 213, 174)
RED = RGBColor(255, 91, 102)
BLUE = RGBColor(82, 155, 255)
WHITE = RGBColor(255, 255, 255)

FONT = "Noto Sans"


def rect(slide, x, y, w, h, fill, radius=True, line=None, transparency=0):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    shape.line.color.rgb = line or fill
    return shape


def textbox(slide, text, x, y, w, h, size=18, color=CREAM, bold=False,
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0,
            font=FONT, line_spacing=1.0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = margin
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def rich_text(slide, runs, x, y, w, h, size=18, align=PP_ALIGN.LEFT,
              valign=MSO_ANCHOR.TOP, margin=0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = margin
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for text, color, bold, run_size in runs:
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(run_size or size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def bullet_list(slide, items, x, y, w, h, size=17, color=MUTED, accent=MINT,
                gap=7, bullet="•"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet}  {item}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(gap)
        p.line_spacing = 1.08
    return box


def add_cover_image(slide, path, x, y, w, h):
    with Image.open(path) as img:
        iw, ih = img.size
    target = w / h
    source = iw / ih
    pic = slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    if source > target:
        visible = target / source
        crop = (1 - visible) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    elif source < target:
        visible = source / target
        crop = (1 - visible) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    return pic


def add_contain_image(slide, path, x, y, w, h):
    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    pw, ph = int(iw * scale), int(ih * scale)
    return slide.shapes.add_picture(str(path), x + (w - pw) / 2, y + (h - ph) / 2,
                                    width=pw, height=ph)


def bg(slide):
    rect(slide, 0, 0, W, H, BG, radius=False)


def kicker(slide, text, x=Inches(.55), y=Inches(.35), color=GOLD):
    textbox(slide, text.upper(), x, y, Inches(7.0), Inches(.28), 10, color, True)


def title(slide, text, y=Inches(.67), size=30, color=CREAM):
    textbox(slide, text, Inches(.55), y, Inches(12.2), Inches(.65), size, color, True)


def footer(slide, n, source=None):
    textbox(slide, f"MEDIMATCH  /  {n:02d}", Inches(.55), Inches(7.13), Inches(2.0), Inches(.18), 8, MUTED, True)
    if source:
        textbox(slide, source, Inches(2.6), Inches(7.08), Inches(10.15), Inches(.28), 7.5, MUTED, False, PP_ALIGN.RIGHT)


def metric(slide, value, label, x, y, w, color=GOLD):
    rect(slide, x, y, w, Inches(1.1), PANEL, True, PANEL_2)
    textbox(slide, value, x + Inches(.18), y + Inches(.12), w - Inches(.36), Inches(.48), 25, color, True)
    textbox(slide, label, x + Inches(.18), y + Inches(.68), w - Inches(.36), Inches(.25), 10, MUTED, False)


def chart_card(slide, img, x, y, w, h, caption, ccolor=GOLD):
    """A real survey screenshot on a white card with a coloured caption strip."""
    rect(slide, x, y, w, h, WHITE, True, PANEL_2)
    add_contain_image(slide, img, x + Inches(.1), y + Inches(.1), w - Inches(.2), h - Inches(.2))
    cap = rect(slide, x, y + h, w, Inches(.34), PANEL, True, ccolor)
    textbox(slide, caption, x + Inches(.12), y + h + Inches(.04), w - Inches(.24), Inches(.26),
            11.5, ccolor, True, PP_ALIGN.CENTER)


def survey_row(slide, items, y=2.55):
    """A numeric-order row of survey charts on white cards, each captioned."""
    n = len(items)
    gap = 0.3
    cw = (12.4 - (n - 1) * gap) / n
    ch = cw / 2.42 + 0.2
    x0 = (13.333 - (n * cw + (n - 1) * gap)) / 2
    for i, (img, cap, c) in enumerate(items):
        x = x0 + i * (cw + gap)
        rect(slide, Inches(x), Inches(y), Inches(cw), Inches(ch), WHITE, True, PANEL_2)
        add_contain_image(slide, img, Inches(x + .1), Inches(y + .1), Inches(cw - .2), Inches(ch - .2))
        rect(slide, Inches(x), Inches(y + ch), Inches(cw), Inches(.36), PANEL, True, c)
        textbox(slide, cap, Inches(x + .08), Inches(y + ch + .05), Inches(cw - .16), Inches(.28),
                11, c, True, PP_ALIGN.CENTER)
    return y + ch + 0.36


def showcase(slide, img, caption, h=5.55, y=1.32):
    """A full app screenshot shown whole (never cropped), centred and large."""
    iw = h * 1.60
    x = (13.333 - iw) / 2
    rect(slide, Inches(x - .04), Inches(y - .04), Inches(iw + .08), Inches(h + .08), PANEL_2, True, PANEL_2)
    add_contain_image(slide, img, Inches(x), Inches(y), Inches(iw), Inches(h))
    textbox(slide, caption, Inches(.7), Inches(y + h + .14), Inches(12.0), Inches(.3),
            12, MUTED, False, PP_ALIGN.CENTER)


prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# 1 — Opening
s = prs.slides.add_slide(blank)
add_cover_image(s, ASSETS / "title.png", 0, 0, W, H)
rect(s, 0, 0, W, H, BG, radius=False, transparency=72)
rect(s, Inches(.58), Inches(.65), Inches(8.2), Inches(5.7), BG, True, BG, transparency=18)
textbox(s, "MEDIMATCH  ·  NATIONAL REDISTRIBUTION COMMAND", Inches(.86), Inches(1.0), Inches(7.5), Inches(.3), 11, GOLD, True)
rich_text(s, [
    ("See surplus.\nDetect need.\n", CREAM, True, 38),
    ("Coordinate impact.", MINT, True, 38),
], Inches(.84), Inches(1.48), Inches(7.65), Inches(2.45), 38)
textbox(s, "Leveraging geospatial intelligence\nfor equitable medical-supply redistribution", Inches(.88), Inches(4.08), Inches(6.9), Inches(.9), 18, MUTED)
rect(s, Inches(.88), Inches(5.18), Inches(5.7), Inches(.03), GOLD, radius=False)
textbox(s, "Lesnar Gitonga  ·  USIU–Africa  ·  23 June 2026", Inches(.88), Inches(5.43), Inches(6.6), Inches(.32), 12, CREAM, True)
textbox(s, "Conference demo · synthetic inventory · no patient records", Inches(.88), Inches(5.85), Inches(6.8), Inches(.28), 10, MUTED)


# 2 — Problem and evidence
s = prs.slides.add_slide(blank); bg(s)
kicker(s, "01  /  The coordination gap")
textbox(s, "The problem is not only scarcity.\nIt is disconnected visibility.", Inches(.55), Inches(.67), Inches(12.2), Inches(.9), 26, CREAM, True)
textbox(s, "Facilities can hold usable surplus while others face urgent shortfalls. MediMatch makes both sides visible in one operating picture.", Inches(.58), Inches(1.62), Inches(12.0), Inches(.4), 14, MUTED)
metric(s, "56.3%", "face weekly or monthly stockouts", Inches(.58), Inches(2.25), Inches(2.85), ORANGE)
metric(s, "82.8%", "rate manual methods ineffective or slow", Inches(3.56), Inches(2.25), Inches(2.85), GOLD)
metric(s, "57.8%", "report wastage of valid supplies", Inches(6.54), Inches(2.25), Inches(2.85), RED)
metric(s, "89.1%", "willing to pilot the system", Inches(9.52), Inches(2.25), Inches(2.85), MINT)
rect(s, Inches(.58), Inches(3.72), Inches(12.0), Inches(2.25), PANEL, True, PANEL_2)
textbox(s, "A supply gap becomes\na coordination question", Inches(.88), Inches(4.0), Inches(5.3), Inches(.72), 19, CREAM, True)
textbox(s, "Which available source is close enough, sufficiently stocked, verified, product-compatible—and practical to move now?", Inches(.88), Inches(4.82), Inches(5.25), Inches(.78), 15, MUTED)
rect(s, Inches(6.65), Inches(4.08), Inches(.04), Inches(1.35), GOLD, radius=False)
textbox(s, "57%", Inches(7.02), Inches(4.02), Inches(1.65), Inches(.62), 34, GOLD, True)
textbox(s, "KEMSA national order-fill rate,\nas of mid-2025 (per the abstract)", Inches(8.35), Inches(4.13), Inches(3.55), Inches(.65), 13, MUTED)
textbox(s, "The baseline is historical; the deck does not present it as the current 2026 rate.", Inches(7.02), Inches(5.08), Inches(4.8), Inches(.45), 10, MUTED)
footer(s, 2, "Sources: MediMatch Nairobi County field study (n=64 healthcare professionals, 80% response, 2025); KEMSA order-fill rate 57%, mid-2025")


# 3 — Survey: respondent profile (Q1–Q3)
s = prs.slides.add_slide(blank); bg(s)
kicker(s, "02  /  The field evidence")
title(s, "Who we surveyed", size=28)
textbox(s, "64 healthcare professionals  ·  Nairobi County  ·  80% response rate  ·  unedited Google Forms output",
        Inches(.55), Inches(1.18), Inches(12.2), Inches(.3), 13, GOLD, True)
textbox(s, "A frontline sample — the clinicians and staff who manage stock every day.",
        Inches(.55), Inches(1.62), Inches(12.2), Inches(.3), 14, MUTED)
SURVEY = ASSETS / "survey"
survey_row(s, [
    (SURVEY / "q1_gender.png", "Q1 · Gender — slight female majority", ORANGE),
    (SURVEY / "q2_age.png",    "Q2 · Age — majority under 40",         MINT),
    (SURVEY / "q3_role.png",   "Q3 · Role — clinical staff majority",  BLUE),
], y=2.5)
footer(s, 3, "Source: MediMatch Nairobi County field study · n=64 · 2025")


# 4 — Survey: the current system (Q4–Q6)
s = prs.slides.add_slide(blank); bg(s)
kicker(s, "03  /  The field evidence")
title(s, "The current system, in their words", size=28)
textbox(s, "Questions 4–6  ·  frequency of stockouts, how surplus is handled, and how they rate today’s methods",
        Inches(.55), Inches(1.18), Inches(12.2), Inches(.3), 13, GOLD, True)
textbox(s, "Frequent stockouts, ad-hoc handling, and locating methods they rate as ineffective or slow.",
        Inches(.55), Inches(1.62), Inches(12.2), Inches(.3), 14, MUTED)
survey_row(s, [
    (SURVEY / "q4_stockouts.png",     "Q4 · 56.3% weekly or monthly stockouts", ORANGE),
    (SURVEY / "q5_surplus.png",       "Q5 · Surplus handled ad-hoc",            GOLD),
    (SURVEY / "q6_effectiveness.png", "Q6 · 82.8% ineffective or slow",         RED),
], y=2.5)
footer(s, 4, "Source: MediMatch survey, n=64 · Q4 and Q6 are the abstract’s 56.3% and 82.8% figures, shown from the raw responses")


# 5 — Survey: challenges, value & readiness (Q7–Q9)
s = prs.slides.add_slide(blank); bg(s)
kicker(s, "04  /  The field evidence")
title(s, "Challenges, value & readiness", size=28)
textbox(s, "Questions 7–9  ·  the main system challenges, the benefits they value, and willingness to pilot",
        Inches(.55), Inches(1.18), Inches(12.2), Inches(.3), 13, GOLD, True)
textbox(s, "Wastage tops the operational challenges — and readiness to pilot MediMatch is high.",
        Inches(.55), Inches(1.62), Inches(12.2), Inches(.3), 14, MUTED)
survey_row(s, [
    (SURVEY / "q7_challenges.png", "Q7 · 57.8% cite wastage of supplies", RED),
    (SURVEY / "q8_benefits.png",   "Q8 · Most-valued benefits",           BLUE),
    (SURVEY / "q9_pilot.png",      "Q9 · 89.1% would pilot (Yes + Maybe)", MINT),
], y=2.5)
footer(s, 5, "Source: MediMatch survey, n=64 · Q7 and Q9 are the abstract’s 57.8% and 89.1% figures, shown from the raw responses")


# 6 — How the engine works (concept, no screenshots)
s = prs.slides.add_slide(blank); bg(s)
kicker(s, "05  /  How it works")
title(s, "How the redistribution engine works", size=27)
textbox(s, "Every listing becomes a geocoded surplus offer or an urgent request, then flows through a transparent pipeline.",
        Inches(.55), Inches(1.22), Inches(12.2), Inches(.3), 14, MUTED)

pipeline = [
    ("DETECT", "Parse each listing into a\ngeocoded surplus or need", ORANGE, .6),
    ("RANK",   "Score candidates urgent-first\non multiple signals",   GOLD, 3.72),
    ("ROUTE",  "Resolve real road geometry\nvia OSRM, arc fallback",   MINT, 6.84),
    ("VERIFY", "Coordinator confirms before\nany real-world transfer", BLUE, 9.96),
]
for i, (h, b, c, xx) in enumerate(pipeline):
    x = Inches(xx)
    rect(s, x, Inches(1.92), Inches(2.75), Inches(1.9), PANEL, True, c)
    rect(s, x, Inches(1.92), Inches(2.75), Inches(.08), c, False)
    textbox(s, h, x + Inches(.22), Inches(2.2), Inches(2.3), Inches(.34), 16, c, True, PP_ALIGN.CENTER)
    textbox(s, b, x + Inches(.22), Inches(2.72), Inches(2.3), Inches(.85), 12, MUTED, False, PP_ALIGN.CENTER)
    if i < 3:
        textbox(s, "→", x + Inches(2.79), Inches(2.5), Inches(.3), Inches(.5), 22, GOLD, True, PP_ALIGN.CENTER)

textbox(s, "Ranking signals", Inches(.6), Inches(4.15), Inches(12.2), Inches(.32), 16, CREAM, True)
signals = [
    ("Proximity", "routed distance", MINT),
    ("Urgency", "clinical priority", RED),
    ("Product fit", "category & cold-chain", GOLD),
    ("Quantity", "full or partial cover", BLUE),
    ("Verification", "trusted source", ORANGE),
]
cw = 2.3; gap = 0.18; x0 = (13.333 - (len(signals) * cw + (len(signals) - 1) * gap)) / 2
for i, (name, desc, c) in enumerate(signals):
    x = x0 + i * (cw + gap)
    rect(s, Inches(x), Inches(4.6), Inches(cw), Inches(.92), PANEL, True, c)
    rect(s, Inches(x), Inches(4.6), Inches(.07), Inches(.92), c, False)
    textbox(s, name, Inches(x + .22), Inches(4.72), Inches(cw - .3), Inches(.3), 13, CREAM, True)
    textbox(s, desc, Inches(x + .22), Inches(5.06), Inches(cw - .3), Inches(.3), 10.5, MUTED, False)

rect(s, Inches(.6), Inches(5.85), Inches(12.13), Inches(.78), PANEL, True, PANEL_2)
textbox(s, "Beyond routing: a demand heatmap surfaces clusters, a one-year impact projection models the gain, and a live Copilot answers from current state — all shown next in the demo.",
        Inches(.9), Inches(6.04), Inches(11.5), Inches(.45), 12.5, MUTED, False, PP_ALIGN.CENTER)
footer(s, 6, "Explainable matching · plain-language situation briefs · human verification before any transfer")


# 7 — Architecture
s = prs.slides.add_slide(blank); bg(s)
kicker(s, "06  /  Architecture and safeguards")
textbox(s, "A production-minded spatial stack\nwith a human decision boundary", Inches(.55), Inches(.66), Inches(12.2), Inches(1.05), 27, CREAM, True, PP_ALIGN.CENTER)

boxes = [
    ("Experience", "React · TypeScript\nMapLibre · three.js", Inches(.6), BLUE),
    ("Coordination API", "Node.js · Express\nmatching · auth · messages", Inches(3.72), ORANGE),
    ("Spatial + routing", "PostgreSQL · PostGIS\nOSRM · Nominatim · CARTO", Inches(6.84), MINT),
    ("Intelligence", "Curated fallback\noptional Claude briefs", Inches(9.96), GOLD),
]
for i, (head, body, xx, c) in enumerate(boxes):
    rect(s, xx, Inches(2.0), Inches(2.75), Inches(2.15), PANEL, True, c)
    rect(s, xx, Inches(2.0), Inches(2.75), Inches(.08), c, False)
    textbox(s, head, xx + Inches(.22), Inches(2.38), Inches(2.3), Inches(.34), 17, CREAM, True, PP_ALIGN.CENTER)
    textbox(s, body, xx + Inches(.22), Inches(2.98), Inches(2.3), Inches(.65), 12, MUTED, False, PP_ALIGN.CENTER)
    if i < len(boxes)-1:
        textbox(s, "→", xx + Inches(2.79), Inches(2.75), Inches(.3), Inches(.5), 22, GOLD, True, PP_ALIGN.CENTER)

rect(s, Inches(.6), Inches(4.65), Inches(12.1), Inches(1.4), PANEL, True, PANEL_2)
textbox(s, "Human verification is the final control", Inches(.9), Inches(4.9), Inches(11.5), Inches(.36), 20, CREAM, True, PP_ALIGN.CENTER)
textbox(s, "No diagnosis · no autonomous transfers · synthetic demo inventory · explicit assumptions · coordinator sign-off", Inches(.9), Inches(5.38), Inches(11.5), Inches(.32), 12, MUTED, False, PP_ALIGN.CENTER)
rect(s, Inches(.9), Inches(5.66), Inches(11.3), Inches(.04), GOLD, False)
footer(s, 7, "Production data layer: PostgreSQL 15 + PostGIS 3")


# 8 — Live demonstration
s = prs.slides.add_slide(blank); bg(s)
kicker(s, "07  /  Live demonstration")
title(s, "See it run — live", size=30)
textbox(s, "Switching to the running platform.", Inches(.55), Inches(1.25), Inches(12.2), Inches(.3), 15, MUTED)
demo = [
    ("National command floor", "Detect, rank and route surplus to urgent need across Kenya"),
    ("Demand heatmap", "Where shortfall concentrates, against hubs and corridors"),
    ("Nairobi research base", "The 64-facility study area at street scale"),
    ("Impact projection", "A transparent one-year model of the gain"),
    ("Copilot", "Live coordination questions, answered from current state"),
]
for i, (h, d) in enumerate(demo):
    yy = 2.05 + i * .92
    rect(s, Inches(.9), Inches(yy), Inches(.6), Inches(.6), PANEL, True, GOLD)
    textbox(s, str(i + 1), Inches(.9), Inches(yy + .08), Inches(.6), Inches(.42), 20, GOLD, True, PP_ALIGN.CENTER)
    textbox(s, h, Inches(1.75), Inches(yy + .02), Inches(6.5), Inches(.34), 18, CREAM, True)
    textbox(s, d, Inches(1.75), Inches(yy + .36), Inches(10.5), Inches(.3), 13, MUTED, False)
footer(s, 8, "Recommended demo uses curated local briefs and synthetic inventory")


# 9 — Close
s = prs.slides.add_slide(blank); bg(s)
kicker(s, "08  /  Closing")
rich_text(s, [
    ("Location is not just a field.\n", CREAM, True, 35),
    ("It is a coordination advantage.", MINT, True, 35),
], Inches(.6), Inches(1.05), Inches(8.0), Inches(1.55), 35)
takeaways = [
    ("Evidence-based", "64 healthcare professionals quantified the coordination gap."),
    ("Location-aware", "Real-time, road-routed matching of surplus to urgent need."),
    ("Explainable", "Transparent, multi-factor scoring — no black box."),
    ("Human-governed", "Coordinator verification before every transfer."),
]
for i, (h, d) in enumerate(takeaways):
    yy = 2.95 + i * .72
    rect(s, Inches(.65), Inches(yy), Inches(.09), Inches(.5), MINT, False)
    textbox(s, h, Inches(.92), Inches(yy - .02), Inches(2.3), Inches(.3), 15, CREAM, True)
    textbox(s, d, Inches(.92), Inches(yy + .26), Inches(5.4), Inches(.3), 12, MUTED, False)
rect(s, Inches(7.1), Inches(2.75), Inches(5.45), Inches(3.0), PANEL, True, GOLD)
textbox(s, "MEDIMATCH", Inches(7.55), Inches(3.22), Inches(4.5), Inches(.45), 27, CREAM, True, PP_ALIGN.CENTER)
textbox(s, "See surplus. Detect need.\nCoordinate impact.", Inches(7.55), Inches(3.95), Inches(4.5), Inches(.9), 20, GOLD, True, PP_ALIGN.CENTER)
textbox(s, "Thank you", Inches(7.55), Inches(5.05), Inches(4.5), Inches(.4), 18, MUTED, False, PP_ALIGN.CENTER)
footer(s, 9, "Dr Lesnar Gitonga · USIU-Africa · Global Public Health 2026")


# 10 — References
def ref_column(slide, items, x, y, w, size=10.5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(5.4))
    tf = box.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = it
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = MUTED
        p.space_after = Pt(9)
        p.line_spacing = 1.04

s = prs.slides.add_slide(blank); bg(s)
kicker(s, "09  /  References")
title(s, "References & sources", size=27)
left_refs = [
    "Davis, F. D. (1989). Perceived usefulness, perceived ease of use, and user acceptance of information technology. MIS Quarterly, 13(3), 319–340.",
    "HealthcareMEA. (2025). Supply chain challenges and the impact of wastage in East African healthcare.",
    "Kenya Medical Supplies Authority (KEMSA). (2025). Supply chain performance and order-fill rates: Mid-year report 2025. Nairobi: KEMSA.",
    "Kenya Ministry of Health. (2023). National health supply chain strategy and universal health coverage. Nairobi: Ministry of Health.",
    "Kenya News Agency. (2025). Addressing stockouts: Inter-county redistribution measures.",
    "Lee, H. L., Padmanabhan, V., & Whang, S. (1997). Information distortion in a supply chain: The bullwhip effect. Management Science, 43(4), 546–558.",
]
right_refs = [
    "MedShare. (2024). Recovering surplus medical supplies: Global impact report.",
    "Ministry of Health Kenya. (2020). Kenya National e-Health Policy 2016–2030. Nairobi: Ministry of Health.",
    "Mugenda, O. M., & Mugenda, A. G. (2003). Research methods: Quantitative and qualitative approaches. Nairobi: Acts Press.",
    "Nation Africa. (2024). Global Fund audit reveals inconsistencies in Kenya’s stock management. Daily Nation.",
    "NHS England. (2022). NHS supply chain redistribution programme: Lessons from the COVID-19 pandemic. London: NHS.",
    "Republic of Kenya. (2019). The Data Protection Act. Nairobi: Government Printer.",
]
ref_column(s, left_refs, .58, 1.55, 5.95)
ref_column(s, right_refs, 6.83, 1.55, 5.95)
footer(s, 10, "Primary data: MediMatch Nairobi County field study (n=64 healthcare professionals, 80% response, 2025)")


prs.core_properties.title = "MediMatch — Geospatial Intelligence for Medical-Supply Redistribution"
prs.core_properties.subject = "Conference presentation"
prs.core_properties.author = "Lesnar Gitonga"
prs.core_properties.keywords = "MediMatch, geospatial, PostGIS, public health, Kenya, medical supplies"
prs.save(OUT)
print(OUT)
